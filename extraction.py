import re
from pathlib import Path
from typing import List
from whisper_service import WhisperService
import fitz  # PyMuPDF

class ExtractionLogic:
    """
    Handles extraction of text from various formats and 
    parsing of Obsidian templates.
    """
    
    @staticmethod
    def extract_pdf_text(pdf_path: Path) -> str:
        """
        Extracts text from a PDF file using PyMuPDF.
        Ensures robust multi-page handling by iterating through all pages.
        """
        text = []
        try:
            with fitz.open(pdf_path) as doc:
                for page in doc:
                    text.append(page.get_text())
            return "\n".join(text)
        except Exception as e:
            return f"Error extracting PDF text: {str(e)}"

    @staticmethod
    def extract_pptx_text(pptx_path: Path) -> str:
        """
        Extracts text from a PPTX file using python-pptx.
        Extracts text from all shapes, text boxes, and speaker notes.
        """
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            text_elements = []

            for slide in prs.slides:
                # Extract text from shapes and text boxes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text_elements.append(shape.text.strip())
                
                # Extract text from speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        text_elements.append(notes.strip())

            return "\n".join(text_elements)
        except Exception as e:
            return f"Error extracting PPTX text: {str(e)}"

    @staticmethod
    def extract_audio_text(audio_path: Path) -> str:
        """
        Extracts text from audio files using the Whisper API.
        """
        service = WhisperService()
        return service.transcribe_audio(audio_path)

    @staticmethod
    def parse_template_headers(template_content: str) -> List[str]:
        """
        Extracts H2 headers from an Obsidian template.
        Example: ## Notes -> 'Notes'
        """
        pattern = r'^##\s+(.+)$'
        headers = re.findall(pattern, template_content, re.MULTILINE)
        return headers

    @staticmethod
    def map_content_to_headers(content: str, headers: List[str]) -> dict:
        """
        Simplistic mapping of content to headers for testing.
        In production, this will be handled by the LLM.
        """
        return {header: f"Content for {header} based on: {content}" for header in headers}
